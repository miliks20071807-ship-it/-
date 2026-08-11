"""
Дизайн-агент, друга команда: /дизайн_презентація <опис ідеї> генерує
текстовий пітч через Claude і рендерить його в pptx — окремо від
/дизайн (HTML-мокап екрана) і /презентація (звіт по задачах команди
з presentation.py, дані з storage.py).

Без нових залежностей: той самий python-pptx, що вже в проєкті для
presentation.py.

Основний рушій — безкоштовний тариф Google Gemini (без карти —
aistudio.google.com/app/apikey), щоб не платити за Claude на кожен
пітч. Якщо GEMINI_API_KEY не задано чи запит впав — автоматичний
фолбек на Claude.
"""

from __future__ import annotations

import json
import os
from pathlib import Path

from anthropic import Anthropic
from pptx import Presentation

import api_usage
from agents.common import call_gemini, extract_text, load_product_description

CLAUDE_MODEL = "claude-sonnet-5"
GEMINI_MODEL = "gemini-3.6-flash"
OUTPUT_PATH = Path(__file__).parent / "pitch.pptx"

SYSTEM_PROMPT = """Ти — продуктовий стратег, що готує короткий пітч нової
фічі чи ідеї для застосунку, описаного нижче.

Продукт:
{product}

Отримуєш короткий опис ідеї й повертаєш структуру пітчу — конкретну,
без води, орієнтовану на невелику продуктову команду (2-5 людей), яка
одразу вирішуватиме, чи братись за це.

Відповідай ЛИШЕ JSON:
{{
  "title": "назва ідеї",
  "subtitle": "одне речення суті",
  "problem": "яку проблему користувача це вирішує",
  "solution": "як саме це працює в застосунку",
  "how_it_works": ["крок або механіка 1", "крок або механіка 2", "..."],
  "next_steps": ["конкретний наступний крок 1", "..."]
}}
"""


def _generate_via_claude(system: str, description: str) -> str:
    """Фолбек на Claude, коли Gemini недоступний. Перед платним викликом
    перевіряє денний ліміт Claude API (api_usage.guard()) — щоб
    недоступність безкоштовного Gemini не призводила до мовчазного
    накопичення рахунку понад ліміт."""
    api_usage.guard()
    client = Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
    response = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=8000,
        system=system,
        thinking={"type": "adaptive"},
        output_config={"effort": "xhigh"},
        messages=[{"role": "user", "content": description}],
    )
    api_usage.record_call()
    return extract_text(response)


def _generate_pitch(description: str) -> dict:
    system = SYSTEM_PROMPT.format(product=load_product_description())

    try:
        raw = call_gemini(system, description, model=GEMINI_MODEL, max_tokens=4000)
    except Exception as e:  # noqa: BLE001 — навмисно широко: будь-яка помилка безкоштовного API веде на фолбек, а не крашить фічу
        print(f"[пітч] Gemini недоступний ({e}), фолбек на Claude")
        raw = _generate_via_claude(system, description)

    raw = raw.strip()
    raw = raw.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
    pitch = json.loads(raw)
    pitch.setdefault("title", "Ідея")
    pitch.setdefault("subtitle", "")
    pitch.setdefault("problem", "")
    pitch.setdefault("solution", "")
    pitch.setdefault("how_it_works", [])
    pitch.setdefault("next_steps", [])
    return pitch


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_text_slide(prs: Presentation, title: str, body_text: str) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = body_text or "—"


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    if not bullets:
        body.text = "—"
        return

    first = True
    for bullet in bullets:
        if first:
            body.text = str(bullet)
            first = False
        else:
            body.add_paragraph().text = str(bullet)


def generate_pitch_deck(description: str) -> Path:
    """Генерує pitch.pptx за описом ідеї і повертає шлях до файлу."""
    pitch = _generate_pitch(description)

    prs = Presentation()
    _add_title_slide(prs, pitch["title"], pitch["subtitle"])
    _add_text_slide(prs, "Проблема", pitch["problem"])
    _add_text_slide(prs, "Рішення", pitch["solution"])
    _add_bullets_slide(prs, "Як це працює", pitch["how_it_works"])
    _add_bullets_slide(prs, "Наступні кроки", pitch["next_steps"])

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH
