"""
Генерація pptx-звіту команди за період (крок 4, команда /презентація).

Дані беруться з того самого storage.py (tasks.json), нових залежностей
для зберігання не додається — лише python-pptx для рендеру файлу.
"""

from __future__ import annotations

from datetime import datetime, timedelta, timezone
from pathlib import Path

from pptx import Presentation

import storage

OUTPUT_PATH = Path(__file__).parent / "report.pptx"


def _parse_iso(value: str) -> datetime:
    return datetime.fromisoformat(value)


def _in_period(task: dict, since: datetime | None) -> bool:
    if since is None:
        return True
    timestamp = task.get("done_at") or task.get("created_at")
    return _parse_iso(timestamp) >= since


def _add_title_slide(prs: Presentation, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Звіт команди"
    slide.placeholders[1].text = subtitle


def _add_summary_slide(prs: Presentation, done_count: int, open_count: int) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Підсумок"
    body = slide.placeholders[1].text_frame
    body.text = f"Закрито задач: {done_count}"
    body.add_paragraph().text = f"Відкрито задач: {open_count}"
    body.add_paragraph().text = f"Разом: {done_count + open_count}"


def _add_task_list_slide(prs: Presentation, title: str, tasks: list[dict]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    if not tasks:
        body.text = "—"
        return

    first = True
    for t in tasks:
        line = f"#{t['id']} {t['text']} ({t['author']})"
        if first:
            body.text = line
            first = False
        else:
            body.add_paragraph().text = line


def build_report(days: int | None = 7) -> Path:
    """Генерує report.pptx за останні `days` днів (None = за весь час)
    і повертає шлях до файлу."""
    since = datetime.now(timezone.utc) - timedelta(days=days) if days else None

    all_tasks = storage.list_tasks(status=None)
    done = [t for t in all_tasks if t["status"] == "done" and _in_period(t, since)]
    open_ = [t for t in all_tasks if t["status"] == "open" and _in_period(t, since)]

    period_label = f"за останні {days} дн." if days else "за весь час"

    prs = Presentation()
    _add_title_slide(prs, f"{period_label} · {datetime.now(timezone.utc):%Y-%m-%d}")
    _add_summary_slide(prs, len(done), len(open_))
    _add_task_list_slide(prs, "Закриті задачі", done)
    _add_task_list_slide(prs, "Відкриті задачі", open_)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH
